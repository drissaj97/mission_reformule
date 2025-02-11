from pptx import Presentation
import openai
import os
import pandas as pd
import streamlit as st
from io import BytesIO

# Configuration OpenAI (Remplacez par votre clé API)
openai.api_key = "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"


def extract_text_from_pptx(file):
    """
    Extrait le texte structuré d'un fichier PowerPoint en identifiant les sections clés.
    """
    prs = Presentation(file)
    slides_data = []
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            slides_data.append("\n".join(slide_text))
    
    return slides_data, prs

def reformulate_text(text):
    """
    Utilise GPT pour reformuler le texte et le rendre plus structuré et homogène.
    """
    prompt = f"""
    Reformule et homogénéise le texte suivant tout en conservant son sens :
    ---
    {text}
    ---
    """
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[{"role": "system", "content": "Tu es un assistant qui améliore la rédaction des documents professionnels."},
                      {"role": "user", "content": prompt}]
        )
        return response["choices"][0]["message"]["content"]
    except Exception as e:
        print(f"Erreur OpenAI: {e}")
        return text  # Retourne le texte original en cas d'échec

def process_pptx(file):
    """
    Traite un fichier PowerPoint : extraction, reformulation et génération d'un nouveau PPT.
    """
    slides_text, prs = extract_text_from_pptx(file)
    
    for slide, text in zip(prs.slides, slides_text):
        reformulated_text = reformulate_text(text)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shape.text = reformulated_text  # Remplace le texte par la version reformulée
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# Interface utilisateur avec Streamlit
st.set_page_config(page_title="Traitement des Fiches de Mission", layout="centered")
st.title("📄 Traitement automatique des fiches de mission")
st.markdown("**Chargez un fichier PPTX et cliquez sur le bouton pour le reformuler automatiquement.**")

uploaded_file = st.file_uploader("📂 Téléchargez votre fichier PPTX", type=["pptx"], help="Glissez-déposez un fichier ou cliquez pour en sélectionner un.")

if uploaded_file is not None:
    if st.button("🚀 Traiter le fichier"):
        with st.spinner("Traitement en cours..."):
            processed_ppt = process_pptx(uploaded_file)
        st.success("✅ Traitement terminé ! Téléchargez votre fichier ci-dessous.")
        st.download_button(label="⬇️ Télécharger le fichier reformulé", data=processed_ppt, file_name="mission_reformule.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
